import json
import openai
import os
from dotenv import load_dotenv
from .models import PracticeTest, Question, Option, Flashcard, FlashcardSet
from thefuzz import fuzz
import io
from PyPDF2 import PdfReader
from django.http import JsonResponse
from .validate_json import ai_prompt
from docx import Document

load_dotenv()  # Load environment variables from .env file

from openai import OpenAI
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))




def calculate_points(task, score):
    if score == 0:  # 0 points awarded for failure
        return 0

    # Difficulty points
    difficulty_map = {
        'easy': 10,
        'medium': 20,
        'hard': 30,
        'ultra-hard': 50
    }
    difficulty_key = (task.difficulty or 'easy').lower()  # default to 'easy'
    difficulty_points = difficulty_map.get(difficulty_key, 10)

    time_points = max(task.duration * 1, 5)  # 1 point per minute
    time_points = min(time_points, 100)

    accuracy_points = score  # keep direct mapping

    # Extra points for number of questions (for practice tests)
    question_points = 0
    if hasattr(task, 'questions'):
        question_points = min(task.questions.count() * 2, 100)  # 2 points per question, max 100

    total_points = round(difficulty_points + time_points + accuracy_points + question_points)

    return total_points


ALLOWED_EXTENSIONS = {"txt", "pdf", "docx", "pptx"}
MAX_FILE_SIZE_MB = 5  # Max file size
MAX_CONTENT_LENGTH = 5000  # Max number of characters to send to AI

def decode_uploaded_file(file) -> str:
    """
    Extract text content from an uploaded file.
    Truncate to MAX_CONTENT_LENGTH characters.
    Supports .txt, .pdf, .docx, .pptx
    """
    # Check file extension
    ext = file.name.split('.')[-1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext}")

    # Check file size
    file.seek(0, io.SEEK_END)
    size_mb = file.tell() / (1024 * 1024)
    file.seek(0)
    if size_mb > MAX_FILE_SIZE_MB:
        raise ValueError(f"File too large: {size_mb:.2f} MB")

    text_content = ""

    try:
        if ext == "txt":
            text_content = file.read().decode("utf-8", errors="ignore")
        elif ext == "pdf":
            reader = PdfReader(file)
            text_content = "\n".join(
                page.extract_text() or "" for page in reader.pages
            )
        elif ext == "docx":
            doc = Document(file)
            text_content = "\n".join(p.text for p in doc.paragraphs)
        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(file)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            text_content = "\n".join(slides_text)

        # Reduce size for AI
        if len(text_content) > MAX_CONTENT_LENGTH:
            text_content = text_content[:MAX_CONTENT_LENGTH] + "\n...[truncated]"

        return text_content.strip()
    except Exception as e:
        raise ValueError(f"Error reading file: {e}")





import json5  # safer JSON parser
import os
from openai import OpenAI

# def generate_activity(prompt, amount, difficulty, activity_type, extra_file_data=None):
#     try:
#         response = client.responses.create(
#             model="gpt-4o",
#             messages=[
#                 {"role": "system", "content": "You are a helpful assistant designed to output JSON."},
#                 {"role": "user", "content": "Tell me about the capital of France, including its name and population."}
#             ],
#             response_format={"type": "json_object"},
#         )

#         ai_content = response.output_text
#         usage = response.usage if hasattr(response, "usage") else {}

def generate_activity(prompt, amount, difficulty, activity_type, extra_file_data=None):
    """
    Generate an AI-powered activity (practice test or flashcards) safely.
    Returns: (dict, usage) OR (None, None) if error.
    """

    try:

        # Combine prompt with optional file content
        full_prompt = prompt.strip()
        if extra_file_data:
            full_prompt += f"\n\nUse the following file content as context:\n{extra_file_data}"

        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": (
                    "You are a helpful assistant that creates educational content in valid JSON format. "
                    "Return only JSON. Avoid extra text."
                )},
                {"role": "user", "content": ai_prompt(full_prompt, amount, difficulty, activity_type, extra_file_data)}
            ],
            response_format={"type": "json_object"},
            max_tokens=3000,
            temperature=0.7,
        )
        
        ai_content = response.choices[0].message.content


        # --- Safer JSON parsing ---
        parsed_content = None
        errors = []
        for parser in [json5.loads, json.loads]:
            try:
                parsed_content = parser(ai_content)
                break
            except Exception as e:
                errors.append(str(e))

        if parsed_content is None:
            print("❌ AI JSON parse errors:", errors)
            print("AI raw output (first 300 chars):", ai_content[:300])
            return None, None

        usage = response.usage if hasattr(response, "usage") else {}
        return parsed_content, usage

    except Exception as e:
        print(f"❌ Error in generate_activity: {e}")
        return None, None
 
    
def save_activity_from_json(json_data, user, activity_type, duration=30):
    """
    Saves AI-generated activities (PracticeTest or FlashcardSet) into the database.
    :param json_data: dict (parsed JSON from AI)
    :param user: CustomUser instance (the owner)
    :param activity_type: 'practice_test' or 'flashcards'
    :return: Created model instance
    """

    if activity_type == "practice_test":
        # Ensure we’re inside the right JSON object
        data = json_data.get("PracticeTest", json_data)

        practice_test = PracticeTest.objects.create(
            title=data.get("title", "Untitled Test"),
            description=data.get("description", ""),
            subject=data.get("subject", "General"),
            duration=data.get("duration", duration),
            difficulty=data.get("difficulty", "Medium"),
            is_public=data.get("is_public", True),
            owner=user
        )

        # Loop through questions
        for q in data.get("questions", []):
            question = Question.objects.create(
                practice_test=practice_test,
                text=q.get("text", ""),
                question_type=q.get("question_type", "mcq"),
                subject=q.get("subject", practice_test.subject),
                answer=q.get("answer", ""),
                explanation=q.get("explanation", "")
            )

            # Add options if MCQ
            if q.get("question_type") == "mcq":
                for opt in q.get("options", []):
                    Option.objects.create(
                        question=question,
                        text=opt.get("text", ""),
                        is_correct=opt.get("is_correct", False)
                    )

        return practice_test

    elif activity_type == "flashcards":
        data = json_data.get("FlashcardSet", json_data)

        flashcard_set = FlashcardSet.objects.create(
            title=data.get("title", "Untitled Flashcards"),
            description=data.get("description", ""),
            subject=data.get("subject", "General"),
            difficulty=data.get("difficulty", "Medium"),
            owner=user
        )

        for card in data.get("flashcards", []):
            Flashcard.objects.create(
                flashcard_set=flashcard_set,
                front=card.get("front", ""),
                back=card.get("back", "")
            )

        return flashcard_set

    else:
        raise ValueError(f"Unsupported activity_type: {activity_type}")

    
def is_similar_answer(correct_answer, user_answer, threshold=80):
    """
    Returns True if the user's answer is sufficiently similar to the correct answer.
    threshold: 0-100, higher means stricter matching
    """
    if not correct_answer or not user_answer:
        return False

    correct_answer = correct_answer.lower().strip()
    user_answer = user_answer.lower().strip()

    similarity = fuzz.ratio(correct_answer, user_answer)
    return similarity >= threshold, similarity

def deduct_credits(usage, user):
    """
    Deduct user credits based on OpenAI token usage.
    """
    # Make sure usage is not None
    if usage is None:
        print("Usage was found to be none, exiting deduct_credits function (located in myapp/utils.py)")
        return

    # Access attributes directly
    prompt_tokens = getattr(usage, "prompt_tokens", 0)
    completion_tokens = getattr(usage, "completion_tokens", 0)
    total_tokens = getattr(usage, "total_tokens", 0)

    user.ai_credits -= total_tokens
    user.save()


# this returns ai_content, automatically deducts credits
def ai_chat_response(prompt, system_content, user, model="gpt-4o-mini", response_format="text", max_tokens=3000, temperature=0.7, stream=False):
    try:

        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_content},
                {"role": "user", "content": prompt},
            ],
            max_tokens=max_tokens,
            temperature=temperature,
            stream=stream,
            response_format={"type": response_format}
        )

        if stream:
            ai_content = ""
            for chunk in response:
                delta = chunk.choices[0].delta.get("content", "")
                ai_content += delta
            usage = {}
        else:
            ai_content = response.choices[0].message.content
            usage = response.usage if hasattr(response, "usage") else {}

        deduct_credits(usage, user)

        return ai_content

    except Exception as e:
        return f"Error: {str(e)}", {}